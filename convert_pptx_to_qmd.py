#!/usr/bin/env python3
"""
convert_pptx_to_qmd.py

Simple converter: PPTX -> Quarto RevealJS (.qmd)

Usage:
  python3 convert_pptx_to_qmd.py /path/to/Picornavirus.pptx

This script will create a `presentation.qmd` next to the input file and an
`images/` folder containing extracted images. It extracts slide-level text
and images and embeds them into the .qmd file.

Dependencies:
  pip install python-pptx

Limitations:
  - Complex slide layouts, embedded charts and SmartArt may not convert cleanly.
  - Text styling is flattened to plain text.
"""
#!/usr/bin/env python3
"""
convert_pptx_to_qmd.py

Simple, robust PPTX -> Quarto RevealJS (.qmd) converter.

Features:
- Extracts all files from ppt/media/ into an images/ folder.
- Extracts shape images and deduplicates by content hash.
- Maps slide relationship rIds to actual media files so videos (mp4) are embedded
  as playable <video> tags instead of placeholder images.
- Prefers the first slide title as the presentation title in YAML front-matter.
- Forces author/email fields to the requested values.

Usage:
  python3 convert_pptx_to_qmd.py "12- Structure and classification of viruses.pptx"

Dependencies:
  pip install python-pptx

"""

import hashlib
import os
import pathlib
import re
import sys
import zipfile
import unicodedata
from typing import Dict, Tuple, List
from PIL import Image
import logging

from pptx import Presentation


# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s', handlers=[
    logging.StreamHandler(sys.stdout)
])


def clean_text_line(s: str) -> str:
    if s is None:
        return ''
    # remove common non-printable artifacts seen in some PPTX exports
    s = s.replace('\u0000', '')
    s = re.sub(r"M-[\^\w\-]+", '', s)
    s = s.replace('\r', '')
    
    # Add spaces around periods followed immediately by capital letters (common in bullet lists)
    s = re.sub(r'\.([A-Z])', r'. \1', s)
    # Add spaces around colons followed by capital letters
    s = re.sub(r':([A-Z])', r': \1', s)
    
    s = s.strip()
    s = re.sub(r"\s+", ' ', s)
    return s


def sentence_case(s: str) -> str:
    s = s.strip()
    if not s:
        return s
    return s[0].upper() + s[1:]


def shorten_bullet(s: str, max_words: int = None) -> str:
    """Extract full text without truncation unless max_words is specified"""
    if max_words is None:
        return s  # Return full text without truncation
    words = s.split()
    if len(words) <= max_words:
        return s
    return ' '.join(words[:max_words]).rstrip(' ,;:') + '...'


def sha1_bytes(b: bytes) -> str:
    return hashlib.sha1(b).hexdigest()


def extract_all_media_from_pptx(pptx_path: pathlib.Path, images_dir: pathlib.Path) -> Tuple[Dict[str, str], Dict[str, str]]:
    """Extract all files under ppt/media into images_dir.
    Returns (name->hash, hash->name).
    """
    name_to_hash = {}
    hash_to_name = {}
    with zipfile.ZipFile(pptx_path, 'r') as z:
        for zi in z.infolist():
            if zi.filename.startswith('ppt/media/'):
                data = z.read(zi.filename)
                name = pathlib.Path(zi.filename).name
                h = sha1_bytes(data)
                # dedupe by hash
                if h in hash_to_name:
                    name_to_hash[name] = h
                    continue
                out_name = name
                out_path = images_dir / out_name
                with open(out_path, 'wb') as f:
                    f.write(data)
                name_to_hash[name] = h
                hash_to_name[h] = out_name
    return name_to_hash, hash_to_name


def build_slide_rel_map(pptx_path: pathlib.Path, nslides: int) -> Dict[int, Dict[str, str]]:
    """Return a map: slide_index -> { rId: filename }

    Parses ppt/slides/_rels/slideN.xml.rels files to map rIds to media file names.
    """
    rels = {}
    with zipfile.ZipFile(pptx_path, 'r') as z:
        for i in range(1, nslides + 1):
            rel_path = f'ppt/slides/_rels/slide{i}.xml.rels'
            rels[i] = {}
            try:
                data = z.read(rel_path).decode('utf-8')
            except KeyError:
                continue
            # find Target="../media/media1.png" and Id="rId1"
            for m in re.finditer(r'Id="(rId[0-9]+)"[^>]*Target="([^"]+)"', data):
                rid = m.group(1)
                target = m.group(2)
                # normalize target basename
                fname = pathlib.Path(target).name
                rels[i][rid] = fname
    return rels


def extract_images_from_shape(shape, images_dir: pathlib.Path, img_counter: int, hash_to_name: Dict[str, str]) -> Tuple[List[str], int]:
    """If the shape has image data, save it to images_dir and return list of filenames saved.
    Deduplicate using hash_to_name mapping. Returns (list_of_saved_names, new_counter)
    """
    saved = []
    try:
        if not hasattr(shape, 'image'):
            return saved, img_counter
        img = shape.image
        blob = img.blob
        h = sha1_bytes(blob)
        if h in hash_to_name:
            saved_name = hash_to_name[h]
            return [saved_name], img_counter
        ext = img.ext
        if not ext:
            ext = 'png'
        name = f'image{img_counter}.{ext}'
        out_path = images_dir / name
        with open(out_path, 'wb') as f:
            f.write(blob)
        hash_to_name[h] = name
        saved.append(name)
        img_counter += 1
    except Exception:
        pass
    return saved, img_counter


def slide_text(slide) -> str:
    """Extract text from slide shapes, preserving structure for bullet points"""
    parts = []
    for shape in slide.shapes:
        try:
            if hasattr(shape, 'text') and shape.text:
                # Check if this is a text frame with paragraphs (bullet points)
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    text_parts = []
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            text_parts.append(clean_text_line(paragraph.text))
                    if text_parts:
                        parts.append('\n'.join(text_parts))
                else:
                    parts.append(clean_text_line(shape.text))
        except Exception:
            continue
    return '\n\n'.join([p for p in parts if p])


def convert(pptx_path: str, out_qmd: str = None):
    logging.debug(f"Starting conversion for: {pptx_path}")
    pptx_path = pathlib.Path(pptx_path)
    if not pptx_path.exists():
        logging.error(f"File not found: {pptx_path}")
        raise FileNotFoundError(pptx_path)

    logging.debug("Checking if the input file exists and is accessible.")
    if not pptx_path.exists():
        logging.error(f"Input file does not exist: {pptx_path}")
        raise FileNotFoundError(f"Input file not found: {pptx_path}")

    logging.debug("Starting to process the PowerPoint file.")
    try:
        prs = Presentation(str(pptx_path))
    except Exception as e:
        logging.error(f"Failed to open the PowerPoint file: {e}")
        raise

    logging.debug("Successfully opened the PowerPoint file.")

    if out_qmd is None:
        out_qmd = pptx_path.with_suffix('.qmd')
    out_qmd = pathlib.Path(out_qmd)

    out_dir = out_qmd.parent
    images_dir = out_dir / 'images'
    images_dir.mkdir(parents=True, exist_ok=True)

    logging.debug(f"Output directory: {out_dir}")
    logging.debug(f"Images directory: {images_dir}")

    # Extract media
    name_to_hash, hash_to_name = extract_all_media_from_pptx(pptx_path, images_dir)
    logging.debug(f"Extracted media: {name_to_hash.keys()}")

    filename_title = pptx_path.stem.replace('_', ' ')
    pres_title = None
    pres_author = None
    try:
        props = prs.core_properties
        if props.title:
            pres_title = props.title
        if props.author:
            pres_author = props.author
    except Exception as e:
        logging.warning(f"Error reading presentation properties: {e}")

    first_slide_title = None
    try:
        first = prs.slides[0]
        if first.shapes.title and first.shapes.title.text:
            first_slide_title = clean_text_line(first.shapes.title.text)
    except Exception as e:
        logging.warning(f"Error reading first slide title: {e}")

    img_counter = 1
    slides_content = []
    slide_rels = build_slide_rel_map(pptx_path, len(prs.slides))

    for i, slide in enumerate(prs.slides, start=1):
        logging.debug(f"Processing slide {i} of {len(prs.slides)}.")
        # If we captured the first slide title for front-matter, don't emit
        # the first slide body (this avoids author/name blocks or cover images
        # being appended as a separate slide).
        if i == 1 and first_slide_title:
            continue

        slide_lines = []
        # text
        text = slide_text(slide)
        seen_text = set()
        if text:
            for para in text.split('\n\n'):
                p = para.strip()
                if not p:
                    continue
                if p in seen_text:
                    continue
                seen_text.add(p)
                slide_lines.append(p)

        # images from shapes
        saved_images = []
        for shape in slide.shapes:
            imgs, img_counter = extract_images_from_shape(shape, images_dir, img_counter, hash_to_name)
            saved_images.extend(imgs)
        seen_images = set()
        for img in saved_images:
            if img in seen_images:
                continue
            seen_images.add(img)
            slide_lines.append(f'![](images/{img}){{ width=70% }}')

        # relationship media
        slide_xml = slide._element.xml
        rids = re.findall(r'r:embed="(rId[0-9]+)"', slide_xml)
        video_exts = {'.mp4', '.mov', '.wmv', '.avi', '.mkv', '.webm', '.mpeg', '.mpg', '.ogv'}
        for rid in rids:
            filename = slide_rels.get(i, {}).get(rid)
            if not filename:
                continue
            _, ext = os.path.splitext(filename.lower())
            if ext in video_exts:
                if (images_dir / filename).exists():
                    stem = os.path.splitext(filename)[0]
                    slide_lines = [ln for ln in slide_lines if not (ln.startswith('![](') and stem in ln)]
                    slide_lines.append(f'<video controls style="max-width:70%"><source src="images/{filename}" type="video/mp4"></video>')
                else:
                    slide_lines.append(f'[Video: images/{filename}](images/{filename})')
            else:
                if (images_dir / filename).exists() and filename not in seen_images:
                    seen_images.add(filename)
                    slide_lines.append(f'![](images/{filename}){{ width=70% }}')

        # cleanup and prefer png
        slide_lines = [ln for ln in slide_lines if not re.match(r'^\s*\d+\s*$', ln)]
        def prefer_png(ln):
            m = re.match(r'!\[]\(images/([^\)]+)\)\{[^}]*\}', ln)
            if not m:
                return ln
            fname = m.group(1)
            stem = os.path.splitext(fname)[0]
            png = images_dir / (stem + '.png')
            if png.exists():
                return f'![](images/{stem}.png){{ width=70% }}'
            return ln
        slide_lines = [prefer_png(ln) for ln in slide_lines]

        # slide title detection
        title_text = None
        try:
            if slide.shapes.title and slide.shapes.title.text:
                raw_title = slide.shapes.title.text
                title_text = clean_text_line(raw_title)
        except Exception:
            title_text = None

        # normalize and map some common non-English headings to English
        if title_text:
            # remove diacritics and normalize for comparison
            norm = unicodedata.normalize('NFKD', title_text).encode('ASCII', 'ignore').decode('ASCII')
            norm = norm.strip().lower()
            title_map = {
                'objectif': 'Objectives',
                'objectifs': 'Objectives',
                "objectif(s)": 'Objectives'
            }
            if norm in title_map:
                title_text = title_map[norm]
            # remove any slide-lines equal to the raw title or normalized variants
            variants = set()
            raw_norm = unicodedata.normalize('NFKD', title_text).encode('ASCII', 'ignore').decode('ASCII').strip()
            variants.add(title_text.strip())
            if 'raw_title' in locals():
                variants.add(clean_text_line(raw_title).strip())
                variants.add(clean_text_line(raw_title).strip().lower())
                variants.add(raw_norm.strip())
            slide_lines = [ln for ln in slide_lines if ln.strip().lower() not in {v.lower() for v in variants}]

        # dedupe
        seen = set()
        deduped = []
        for ln in slide_lines:
            key = ln.strip()
            if key in seen:
                continue
            seen.add(key)
            deduped.append(ln)
        slide_lines = deduped

        # make bullets - improved processing for better bullet point separation
        bullets = []
        other = []
        for ln in slide_lines:
            s = ln.strip()
            if not s:
                continue
            if s.startswith('![](') or s.startswith('<video') or s.startswith('<source') or s.startswith('['):
                other.append(ln)
                continue

            # Split text that contains multiple sentences or bullet points
            # Look for patterns that suggest multiple bullet points in one text block
            text_parts = []
            if '\n' in s:
                # Text already has line breaks - treat each line as separate bullet
                text_parts = [line.strip() for line in s.split('\n') if line.strip()]
            else:
                # Try to split on sentence endings followed by capital letters (new bullet points)
                # Also split on patterns like ". Word" or "? Word" or "! Word"
                # Additionally split on common bullet point indicators
                sentences = re.split(r'(?<=[.!?])\s+(?=[A-Z])', s)
                if len(sentences) <= 1:
                    # Try other patterns for splitting bullet points
                    # Split on patterns that suggest list items
                    alt_patterns = [
                        r'\.\s+[A-Z]',  # Period followed by capital letter
                        r':\s+[A-Z]',   # Colon followed by capital letter  
                        r';\s+[A-Z]',   # Semicolon followed by capital letter
                        r'\s{2,}[A-Z]', # Multiple spaces followed by capital letter
                    ]
                    for pattern in alt_patterns:
                        parts = re.split(pattern, s)
                        if len(parts) > 1:
                            sentences = parts
                            break
                
                if len(sentences) > 1:
                    text_parts = sentences
                else:
                    text_parts = [s]
            
            # Process each part as a separate bullet point
            for part in text_parts:
                p = clean_text_line(part)
                p = sentence_case(p)
                p = shorten_bullet(p)  # No max_words limit - extract full content
                if p:  # Only add non-empty bullets
                    bullets.append('- ' + p)
        # Remove the truncation of bullets - extract all content

        # assemble
        assembled = []
        # omit repeating the main title on slide 1 if we used first_slide_title as front-matter
        if i == 1 and first_slide_title:
            # include author/subtitle if wanted
            if pres_author:
                assembled.append(f'### {clean_text_line(pres_author)}')
        else:
            if title_text:
                assembled.append(f'## {title_text}')
        if bullets:
            assembled.extend(bullets)
        assembled.extend(other)

        # if a slide produced no usable text/media, skip it to avoid blank slides
        if not assembled or all((ln.strip() == '' for ln in assembled)):
            continue

        slides_content.append('\n\n'.join(assembled))

    # write QMD
    logging.debug(f"Attempting to write the output QMD file to: {out_qmd}")
    try:
        with open(out_qmd, 'w', encoding='utf-8') as f:
            yaml_content = generate_yaml_front_matter(first_slide_title or pres_title or filename_title)
            f.write(yaml_content)
            for slide_md in slides_content:
                f.write(slide_md + '\n\n')
    except Exception as e:
        logging.error(f"Failed to write the output QMD file: {e}")
        raise

def generate_yaml_front_matter(title: str) -> str:
    """Generate the YAML front matter for the QMD file."""
    return f"""---
title: \"{title}\"
authors:
  - name: \"Sopheap Oeng\"
    email: \"sopheap15@gmail.com\"
format:
  revealjs:
    theme: simple
    transition: fade
    width: 1280
    height: 720
    margin: 0.1
    minScale: 0.2
    maxScale: 1.5
    slideNumber: true
---\n\n"""

def convert_emf_to_png(emf_path: pathlib.Path, png_path: pathlib.Path):
    """Convert an EMF file to PNG format."""
    try:
        with Image.open(emf_path) as img:
            img.save(png_path, format="PNG")
    except Exception as e:
        print(f"Failed to convert {emf_path} to PNG: {e}", file=sys.stderr)


def get_unique_filename(directory: pathlib.Path, base_name: str, extension: str) -> pathlib.Path:
    """Generate a unique filename by appending a number if the file exists."""
    counter = 1
    new_name = f"{base_name}.{extension}"
    new_path = directory / new_name
    while new_path.exists():
        new_name = f"{base_name}_{counter}.{extension}"
        new_path = directory / new_name
        counter += 1
    return new_path


def extract_and_convert_emf_to_png(images_dir: pathlib.Path):
    """Find and convert all EMF files in the images directory to PNG."""
    for emf_file in images_dir.glob("*.emf"):
        base_name = emf_file.stem
        png_path = get_unique_filename(images_dir, base_name, "png")
        convert_emf_to_png(emf_file, png_path)
        emf_file.unlink()  # Remove the original EMF file


# Add this call after extracting media in the convert function
    extract_and_convert_emf_to_png(images_dir)

    logging.debug("Script execution completed. Check the output directory for results.")

if __name__ == "__main__":
    if len(sys.argv) < 2:
        logging.error("Usage: python convert_pptx_to_qmd.py <path_to_pptx_file>")
        sys.exit(1)

    pptx_file = sys.argv[1]
    try:
        convert(pptx_file)
        logging.info("Conversion completed successfully.")
    except Exception as e:
        logging.error(f"An error occurred during conversion: {e}")

